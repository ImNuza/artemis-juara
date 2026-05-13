"""
Build the Artemis pitch deck as a .pptx file with 13 formatted slides.
Output: docs/Artemis_Pitch_Deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
OUT_PATH = ROOT / "docs" / "Artemis_Pitch_Deck.pptx"

# -- Colour palette ------------------------------------------------------------
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOLD       = RGBColor(0xF0, 0xB9, 0x3B)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
GREY       = RGBColor(0xBD, 0xBD, 0xBD)
LIGHT_TEXT = RGBColor(0xEC, 0xF0, 0xF1)
DARK_TEXT  = RGBColor(0x2C, 0x3E, 0x50)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

SLIDE_NUM = [0]  # mutable counter so helpers can auto-increment


# -- Helpers -------------------------------------------------------------------
def dark_slide():
    """Add a blank slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    SLIDE_NUM[0] += 1
    return slide


def light_slide():
    """Add a blank slide with white background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_WHITE
    SLIDE_NUM[0] += 1
    return slide


def add_textbox(slide, left, top, width, height, text="", font_size=18,
                color=LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(line_spacing * font_size - font_size)
    return tf


def add_bullets(tf, items, font_size=16, color=LIGHT_TEXT, bold_first=False):
    """Add bullet paragraphs to an existing text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
        if bold_first and i == 0:
            p.font.bold = True
    return tf


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def gold_line(slide, left, top, width, height=None):
    """Thin gold accent line."""
    h = Pt(height) if height else Pt(3)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_footer(slide, slide_num=None):
    """Add a subtle footer bar with team name, competition, and slide number."""
    # Thin gold separator line
    gold_line(slide, 0.8, 7.1, 11.7, 1.5)

    # Footer text
    footer_text = "ImNuza & Xynerss  |  Artemis Quant Competition 2026  |  Track 1"
    add_textbox(slide, 0.8, 7.15, 9.0, 0.3,
                footer_text,
                font_size=9, color=GREY, alignment=PP_ALIGN.LEFT)

    if slide_num is not None:
        add_textbox(slide, 10.5, 7.15, 2.0, 0.3,
                    str(slide_num),
                    font_size=9, color=GREY, alignment=PP_ALIGN.RIGHT)


def add_table(slide, left, top, col_widths, headers, rows,
              header_bg=GOLD, header_fg=BG_DARK, body_bg=None, body_fg=LIGHT_TEXT,
              font_size=14):
    """Add a formatted table. Returns the table shape."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Inches(left), Inches(top),
                                         Inches(total_w), Inches(0.4 * n_rows))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = header_fg
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = body_fg if body_fg else LIGHT_TEXT
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            if body_bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = body_bg
            else:
                cell.fill.solid()
                if ri % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x3A)
                else:
                    cell.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x33)

    return table_shape


def placeholder_box(slide, left, top, width, height, label):
    """Add a dashed-border placeholder for a chart or image."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x3A)
    shape.line.color.rgb = RGBColor(0x3A, 0x3A, 0x50)
    shape.line.width = Pt(1)
    add_textbox(slide, left + 0.3, top + height/2 - 0.3, width - 0.6, 0.6,
                label,
                font_size=13, color=GREY, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 1: TITLE
# =============================================================================
s = dark_slide()
add_textbox(s, 1.5, 1.5, 10.3, 1.5,
            "BTC Regime-Gated Alt Factor Strategy",
            font_size=44, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 3.2, 10.3, 0.8,
            "A systematic approach to crypto investing using Bitcoin as the on/off switch",
            font_size=22, color=GREY, alignment=PP_ALIGN.CENTER)
gold_line(s, 4.5, 4.2, 4.3)
add_textbox(s, 1.5, 4.6, 10.3, 0.6,
            "ImNuza & Xynerss  |  Artemis Quant Competition 2026  |  Track 1",
            font_size=18, color=GREY, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 5.3, 10.3, 0.5,
            "June 2026",
            font_size=16, color=GREY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(s, "Welcome. We are presenting a BTC regime-gated alt factor strategy for Track 1 of the Artemis Quant Competition.")


# =============================================================================
# SLIDE 2: THE PROBLEM
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "The Problem",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)
tf = add_textbox(s, 0.8, 1.8, 5.5, 2.0,
                 "Crypto has massive upside.\nIt also has -64% drawdowns.",
                 font_size=24, bold=True, color=WHITE)
add_bullets(tf, [
    "Altcoins are levered bets on Bitcoin: when BTC drops, alts get crushed.",
    "A portfolio that is always invested takes every drawdown.",
    "Most strategies focus on what to buy. Few ask when to buy.",
], font_size=18, color=LIGHT_TEXT)
add_textbox(s, 0.8, 4.5, 5.5, 0.8,
            "Our question:",
            font_size=22, bold=True, color=GOLD)
add_textbox(s, 0.8, 5.2, 5.5, 1.2,
            "Can we capture altcoin upside while\navoiding the worst drawdowns?",
            font_size=26, bold=True, color=WHITE)
placeholder_box(s, 7.0, 1.5, 5.5, 5.2,
                "[ BTC price chart with -64% drawdown\n  highlighted during 2022 ]")
add_footer(s, 2)
add_speaker_notes(s, "Crypto offers asymmetric upside but comes with punishing drawdowns. BTC dropped 64% in 2022. Alts dropped further. Any strategy that is always long absorbs every one of those drawdowns. Our question is simple: can we capture the upside while dodging the worst pain?")


# =============================================================================
# SLIDE 3: TWO-LAYER ARCHITECTURE
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "Our Answer: Two Layers, One Decision Per Week",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)

# Layer 1 box
shape1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.2))
shape1.fill.solid()
shape1.fill.fore_color.rgb = RGBColor(0x1E, 0x3D, 0x59)
shape1.line.color.rgb = GOLD
shape1.line.width = Pt(2)

add_textbox(s, 1.2, 1.9, 10.8, 0.5,
            "LAYER 1: BTC REGIME GATE",
            font_size=22, bold=True, color=GOLD)
add_textbox(s, 1.2, 2.5, 10.8, 0.4,
            "5 signals, 2yr min-max normalised, weighted composite score (0-100)",
            font_size=16, color=LIGHT_TEXT)
add_textbox(s, 1.2, 3.0, 3.5, 0.5,
            ">= 60  BULL  :  deploy alts", font_size=16, bold=True, color=GREEN)
add_textbox(s, 4.7, 3.0, 3.5, 0.5,
            "<= 35  BEAR  :  flat, protect capital", font_size=16, bold=True, color=RED)
add_textbox(s, 8.2, 3.0, 4.0, 0.5,
            "36-59  NEUTRAL  :  flat, wait", font_size=16, bold=True, color=GREY)

# Arrow
add_textbox(s, 6.0, 4.1, 1.3, 0.5, "v  BULL only",
            font_size=16, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

# Layer 2 box
shape2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.0))
shape2.fill.solid()
shape2.fill.fore_color.rgb = RGBColor(0x1E, 0x47, 0x2E)
shape2.line.color.rgb = GOLD
shape2.line.width = Pt(2)

add_textbox(s, 1.2, 4.8, 10.8, 0.5,
            "LAYER 2: ALT FACTOR ENGINE",
            font_size=22, bold=True, color=GOLD)
add_textbox(s, 1.2, 5.3, 10.8, 0.4,
            "7 Hyperliquid perps ranked by: Funding rate, inverted (45%) + Price momentum, 7-week (55%)",
            font_size=16, color=LIGHT_TEXT)
add_textbox(s, 1.2, 5.8, 10.8, 0.5,
            "Top 3, signal-weighted, 40% single-asset cap  |  Leverage: 1.5x (BTC 60-70) / 2.5x (BTC >= 70)",
            font_size=16, color=LIGHT_TEXT)

add_textbox(s, 0.8, 6.85, 11.7, 0.3,
            "T-1 lag on every signal. Weekly rebalance. 17% time-in-market.",
            font_size=14, color=GREY, alignment=PP_ALIGN.CENTER)
add_footer(s, 3)
add_speaker_notes(s, "Two layers, one decision per week. Layer 1 is a blunt filter: it is not a market-timing signal, it just asks whether the environment is favorable for taking alt risk. Layer 2 is the alpha engine: it ranks assets within the BULL window. The key insight: the gate is a deliberate on/off switch, not a subtle timing signal.")


# =============================================================================
# SLIDE 4: LAYER 1: BTC REGIME GATE
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "Layer 1: The BTC Regime Gate",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)

add_table(s, 0.8, 1.8,
          [5.0, 2.5, 4.2],
          ["Signal", "Weight", "What it measures"],
          [
              ["MVRV Z-Score", "30%", "Is BTC cheap or expensive? (market vs realised cap)"],
              ["Puell Multiple", "25%", "Miner profitability: exhausted sellers at bottoms"],
              ["200-Week MA Band", "20%", "Structural trend: above or below long-term mean"],
              ["Stablecoin Supply (30d)", "15%", "Capital inflows: is dry powder accumulating?"],
              ["BTC/ETH Dominance (30d ROC)", "10%", "Risk appetite: is money rotating toward risk?"],
          ],
          font_size=14)

add_textbox(s, 0.8, 5.2, 5.5, 1.5,
            "Every signal: rolling 2-year min-max normalised.\n"
            "All data lagged by 1 week; no peeking at the future.\n\n"
            "The gate kept us flat through all of 2022.\n"
            "BTC -64%  |  Strategy 0%.",
            font_size=17, color=LIGHT_TEXT)

placeholder_box(s, 7.0, 5.0, 5.5, 2.0,
                "[ BTC regime bands chart:\n  BULL green / BEAR red / NEUTRAL grey\n  overlaid on BTC price, 2022-2026 ]")
add_footer(s, 4)
add_speaker_notes(s, "Five signals, each economically motivated. Rolling 2-year normalization prevents lookahead bias. The 2022 result is the cleanest illustration: BTC dropped 64%, our strategy was flat at 0%. The gate did its job.")


# =============================================================================
# SLIDE 5: LAYER 2: ALT FACTOR ENGINE
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "Layer 2: The Alt Factor Engine",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)

add_table(s, 0.8, 1.8,
          [3.5, 2.0, 6.2],
          ["Factor", "Weight", "Logic"],
          [
              ["Funding Rate (inverted)", "45%",
               "Negative funding = shorts paying premium = bullish. Cross-sectional rank."],
              ["Price Momentum (7-week)", "55%",
               "7-week trailing return, cross-sectionally ranked. Dominant factor."],
          ],
          font_size=15)

add_textbox(s, 0.8, 3.6, 5.5, 0.5,
            "Portfolio rules (BULL weeks only):",
            font_size=20, bold=True, color=GOLD)

tf = add_textbox(s, 0.8, 4.2, 5.5, 2.5, "", font_size=17, color=LIGHT_TEXT)
add_bullets(tf, [
    "Rank 7 assets by composite score",
    "Select top 3, weight by signal strength",
    "40% single-asset cap (tested; positive expected value)",
    "BTC >= 70 : 2.5x  |  BTC 60-70 : 1.5x",
], font_size=17, color=LIGHT_TEXT)

# Universe
add_textbox(s, 7.5, 3.6, 5.0, 0.5,
            "Investable Universe (7 assets)",
            font_size=20, bold=True, color=GOLD)
add_table(s, 7.5, 4.2,
          [2.2, 2.8],
          ["Category", "Assets"],
          [
              ["L1 Tokens", "SOL, HYPE"],
              ["Exchange Token", "BNB"],
              ["Privacy / Diversifier", "XMR"],
              ["Crypto Equities", "COIN, HOOD"],
              ["Stablecoin Equity", "CRCL"],
          ],
          font_size=13)

add_textbox(s, 7.5, 6.65, 5.0, 0.3,
            "Cross-category avg correlation: 0.06",
            font_size=14, color=GREY, alignment=PP_ALIGN.CENTER)
add_footer(s, 5)
add_speaker_notes(s, "Two factors survived from five tested. Momentum is dominant at -0.35 Sharpe delta when removed. Funding adds diversification: removing it widens drawdown from -27% to -35%. The 7 assets span five categories with cross-category correlation of just 0.06; the ranking model has real differentiation to work with.")


# =============================================================================
# SLIDE 6: HEADLINE RESULTS
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "Headline Results  (Jan 2022 - May 2026)",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)

add_table(s, 0.8, 1.8,
          [3.5, 2.8, 2.8, 3.0],
          ["Metric", "Our Strategy", "BTC B&H", "EW Alts (no gate)"],
          [
              ["Annual return", "48.8%", "16.3%", "40.7%"],
              ["Sharpe ratio", "1.31", "0.32", "0.81"],
              ["Max drawdown", "-27.2%", "-64.3%", "-66.6%"],
              ["Calmar ratio", "1.79", "0.25", "0.61"],
              ["Final equity", "5.63x", "1.93x", "4.41x"],
              ["Time in market", "17%", "100%", "100%"],
              ["Win rate (traded)", "71.1%", "50.4%", "57.1%"],
          ],
          font_size=16)

add_textbox(s, 0.8, 5.8, 11.7, 0.5,
            "Net of Hyperliquid fees and slippage: Sharpe 1.27, 5.33x  |  Cost drag: 3.1% of Sharpe",
            font_size=16, color=GREY, alignment=PP_ALIGN.CENTER)

tf = add_textbox(s, 0.8, 6.5, 11.7, 0.6, "",
                 font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
tf.paragraphs[0].text = "We beat both benchmarks on every metric; invested only 17% of the time."
add_footer(s, 6)
add_speaker_notes(s, "Headline numbers. Sharpe 1.31, max DD -27.2%, 5.63x final equity. Net of costs: Sharpe 1.27, 5.33x. We beat BTC buy-and-hold and the equal-weight alt portfolio on every metric. The Calmar ratio tells the cleanest story: 1.79 vs 0.25 for BTC; that is the regime gate earning its keep.")


# =============================================================================
# SLIDE 7: WHERE RETURNS COME FROM
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "Where the Returns Come From",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)

add_table(s, 0.8, 1.8,
          [2.2, 2.2, 2.2, 2.2, 3.2],
          ["Year", "Return", "Sharpe", "BULL weeks", "Context"],
          [
              ["2022", "0.0%", "-", "0", "Flat all year: gate did its job"],
              ["2023", "+41.7%", "1.20", "2", "Modest deployment, strong return"],
              ["2024", "+205.0%", "3.17", "24", "Dominant year: bulk of BULL weeks"],
              ["2025", "+42.7%", "1.77", "12", "Second positive cycle observation"],
              ["2026 YTD", "0.0%", "-", "0", "Flat: BEAR/NEUTRAL since Nov 2025"],
          ],
          font_size=15)

add_textbox(s, 0.8, 5.2, 11.7, 0.5,
            "Excluding 2024: Sharpe 1.18, 2.26x equity across the remaining 175 weeks.",
            font_size=18, bold=True, color=WHITE)

tf = add_textbox(s, 0.8, 5.9, 11.7, 1.0, "", font_size=17, color=LIGHT_TEXT)
add_bullets(tf, [
    "2024 is unambiguously the dominant year. We do not minimise this.",
    "The framework still produces positive risk-adjusted returns across multiple sub-periods.",
    "Whether 2024 repeats is unknowable. The most we can claim: the framework works across regimes.",
], font_size=17, color=LIGHT_TEXT)
add_footer(s, 7)
add_speaker_notes(s, "Year-by-year breakdown. 2024 dominates with 24 of 38 lifetime BULL weeks. But the framework produces positive risk-adjusted returns in 2023 and 2025 as well. Excluding 2024 entirely, the strategy still posts Sharpe 1.18 across the remaining 175 weeks. We are honest about the concentration.")


# =============================================================================
# SLIDE 8: FACTOR ATTRIBUTION
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "What Drives Performance",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)

add_table(s, 0.8, 1.8,
          [3.5, 2.5, 2.8, 2.8],
          ["Factor removed", "Sharpe", "Max DD", "Delta Sharpe"],
          [
              ["None (baseline)", "1.31", "-27.2%", "-"],
              ["Price momentum", "0.96", "-23.3%", "-0.35"],
              ["Funding rate", "1.18", "-35.4%", "-0.13"],
          ],
          font_size=17)

add_textbox(s, 0.8, 3.8, 5.5, 1.5,
            "Momentum carries the signal.\nFunding adds diversification.\n\n"
            "Momentum-only: Sharpe 1.18, DD -35.4%\n"
            "Funding-only:  Sharpe 0.96, DD -23.3%\n"
            "Combined:       Sharpe 1.31, DD -27.2%",
            font_size=17, color=LIGHT_TEXT)

add_textbox(s, 7.5, 3.8, 5.0, 2.0,
            "112 configurations tested.\n\n"
            "Shipped config is NOT the in-sample\n"
            "Sharpe maximum (1.46 at threshold 55).\n\n"
            "We chose threshold 60 for drawdown\n"
            "control and walk-forward conservatism.\n\n"
            "We are explicitly not running at\n"
            "the in-sample Sharpe peak.",
            font_size=16, color=LIGHT_TEXT)

placeholder_box(s, 0.8, 5.6, 5.5, 1.3,
                "[ Factor attribution bar chart ]")
add_footer(s, 8)
add_speaker_notes(s, "Factor attribution. Momentum is the dominant factor: removing it costs 0.35 Sharpe. Funding adds 0.13 Sharpe and keeps drawdown narrower. We tested 112 configurations; the shipped values were chosen for risk-adjusted conservatism, not maximum in-sample Sharpe.")


# =============================================================================
# SLIDE 9: ROBUSTNESS
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "Robustness: Walk-Forward and Monte Carlo",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)

# Walk-forward
add_textbox(s, 0.8, 1.8, 5.5, 0.5, "Walk-Forward Validation",
            font_size=22, bold=True, color=GOLD)
add_table(s, 0.8, 2.4,
          [2.5, 1.5, 1.5, 1.5],
          ["Window", "Sharpe", "Return", "BULL weeks"],
          [
              ["Train (2022-2023)", "1.25", "34.8%", "4 of 104"],
              ["Test (2024-2026)", "2.07", "112.1%", "65 of 122"],
          ],
          font_size=14)

tf = add_textbox(s, 0.8, 3.9, 5.5, 1.5, "", font_size=15, color=LIGHT_TEXT)
add_bullets(tf, [
    "Threshold optimised on low-BULL window.",
    "Test Sharpe positive, driven by the bull market.",
    "We ship 60; more conservative than train optimum of 50.",
], font_size=15, color=LIGHT_TEXT)

# Monte Carlo
add_textbox(s, 7.5, 1.8, 5.0, 0.5, "Monte Carlo (250 runs)",
            font_size=22, bold=True, color=GOLD)
add_table(s, 7.5, 2.4,
          [2.5, 1.5],
          ["Percentile", "Sharpe"],
          [
              ["5th", "0.10"],
              ["25th", "0.48"],
              ["50th (median)", "0.77"],
              ["75th", "1.02"],
              ["95th", "1.23"],
              ["Shipped", "1.31"],
          ],
          font_size=14)

tf = add_textbox(s, 7.5, 4.8, 5.0, 2.0, "", font_size=15, color=LIGHT_TEXT)
add_bullets(tf, [
    "24.8% of runs exceed Sharpe 1.0.",
    "82.0% exceed Sharpe 0.5.",
    "Shipped sits in upper tail: honest about specification sensitivity.",
    "Profitable across most configs, but the chosen config captures the upper end.",
], font_size=15, color=LIGHT_TEXT)

add_textbox(s, 0.8, 5.8, 11.7, 0.5,
            "[ Walk-forward equity curves  +  Monte Carlo Sharpe histogram ]",
            font_size=14, color=GREY, alignment=PP_ALIGN.CENTER)
add_footer(s, 9)
add_speaker_notes(s, "Robustness checks. Walk-forward confirms directional generalization but the test window is structurally easier (65 of 104 lifetime BULL weeks fall in test). Monte Carlo with 250 randomized runs shows the strategy is profitable across most configurations, but our shipped value sits in the upper tail. We flag this, not hide it.")


# =============================================================================
# SLIDE 10: TESTED AND REJECTED
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "What We Tested and Rejected",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)

add_textbox(s, 0.8, 1.6, 11.7, 0.5,
            "12 ideas tested. 2 survived. Every rejection documented with dates and metrics.",
            font_size=18, bold=True, color=WHITE)

add_table(s, 0.8, 2.3,
          [4.5, 2.5, 4.7],
          ["Test", "Result", "Disposition"],
          [
              ["4 extra factors (revenue, activity, OI, ETF)", "All +0.01 Sharpe or no data", "Removed"],
              ["Short overlay during BEAR markets", "Sharpe 0.05 to 0.48, DD -75% to -86%", "Rejected: destroyed capital"],
              ["BTC as 8th tradeable asset", "Sharpe 1.13 (-0.18), 4.41x (-1.22x)", "Rejected: gate beats position"],
              ["Continuous leverage scaling", "No improvement", "Rejected: discrete is cleaner"],
              ["Adaptive vol-scaled thresholds", "Degraded Sharpe", "Rejected: fixed more robust"],
              ["TAO shadow asset", "Date-alignment bug invalidated test", "Removed"],
          ],
          font_size=13)

add_textbox(s, 0.8, 5.8, 11.7, 0.8,
            "Short overlay tested twice against different baselines. Destroyed capital both times.\n"
            "The factor model ranks for longs; that does not invert to ranking for shorts.",
            font_size=18, bold=True, color=RED)
add_footer(s, 10)
add_speaker_notes(s, "This slide is critical for the Critical Evaluation score. We tested 12 ideas. Only 2 survived. The short overlay is the most important rejection: tested twice, destroyed capital both times. The factor model identifies hated, oversold alts for longs; in bear-market rallies, those exact names bounce hardest. Long rank does not equal short rank. Every rejection is in the report with dates and specific metrics.")


# =============================================================================
# SLIDE 11: HONEST LIMITATIONS
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "Honest Limitations",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)

add_textbox(s, 0.8, 1.6, 11.7, 0.5,
            "What could break this strategy; and what we are not claiming.",
            font_size=18, bold=True, color=WHITE)

limitations = [
    ("1", "Survivorship bias",
     "All 7 assets survived 2022-2026 and most performed strongly. A backtest universe chosen with hindsight overstates live performance. This is the dominant risk."),
    ("2", "One-cycle validation",
     "About 4.25 years, one crypto cycle. Bounded by Hyperliquid's history. Factor relationships may shift in structurally different regimes."),
    ("3", "Funding mechanism change",
     "The funding-rate factor depends on HL's 8-hour cadence and 0.01% cap. A change alters the signal with no historical precedent."),
    ("4", "BTC-alt decorrelation",
     "If alts rally while BTC is bearish, the gate keeps us flat through gains we should have captured."),
    ("5", "Equity perp regulatory risk",
     "COIN, HOOD, CRCL on HL HIP-3. Regulatory action shrinks the universe from 7 to 4 assets."),
    ("6", "Specification sensitivity",
     "Six tuned parameters. Monte Carlo median Sharpe 0.77 vs our 1.31. The gap is real and we flag it."),
]
y = 2.3
for num, title, desc in limitations:
    add_textbox(s, 0.8, y, 0.4, 0.4, num,
                font_size=20, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_textbox(s, 1.3, y, 2.5, 0.4, title,
                font_size=17, bold=True, color=WHITE)
    add_textbox(s, 3.9, y, 8.6, 0.7, desc,
                font_size=14, color=LIGHT_TEXT)
    y += 0.8

add_textbox(s, 0.8, 6.95, 11.7, 0.2,
            "We are not claiming this strategy will return 49% annually going forward. "
            "We are claiming the framework is sound, the thinking is honest, and the edge is economically motivated.",
            font_size=15, color=GREY, alignment=PP_ALIGN.CENTER)
add_footer(s, 11)
add_speaker_notes(s, "Six honest limitations. Survivorship bias is number one: every asset performed well. One-cycle validation: we do not know how this holds in different regimes. We are not claiming 49% forward returns. We are claiming the framework is sound and the thinking is honest.")


# =============================================================================
# SLIDE 12: WHY THIS WINS
# =============================================================================
s = dark_slide()
add_textbox(s, 0.8, 0.5, 11.7, 0.8, "Why This Wins",
            font_size=36, bold=True, color=GOLD)
gold_line(s, 0.8, 1.2, 3.0)

add_textbox(s, 0.8, 1.5, 11.7, 0.5,
            "Not the highest Sharpe in the competition. The best thinking.",
            font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Four judging criteria as cards
boxes = [
    ("Research Quality", "30%",
     ["T-1 lag on every signal",
      "Rolling 2yr normalisation (not full-sample)",
      "All data sources documented",
      "Reproducible pipeline (GitHub + Quick Start)",
      "Funding costs deducted from PnL"]),
    ("Signal / Edge Validity", "30%",
     ["Regime gate economically motivated",
      "Funding rate: crypto-native, no TradFi equivalent",
      "Each factor mapped to academic or practitioner precedent",
      "Factor attribution quantifies marginal contribution"]),
    ("Critical Evaluation", "20%",
     ["12 tested-and-rejected paths with dates",
      "Survivorship bias called out as dominant risk",
      "Walk-forward and Monte Carlo with honest framing",
      "We are not running at the in-sample Sharpe peak"]),
    ("Communication", "20%",
     ["Two-layer architecture, one decision per week",
      "Every number has a story",
      "Every limitation has a 'why'",
      "Clear architecture diagram and regime visualisation"]),
]

x_positions = [0.8, 4.0, 7.2, 10.4]
for i, (title, weight, items) in enumerate(boxes):
    x = x_positions[i]
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(2.2), Inches(2.8), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x3A)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1)

    add_textbox(s, x + 0.2, 2.35, 2.4, 0.4, title,
                font_size=16, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.2, 2.7, 2.4, 0.3, weight,
                font_size=14, color=GREY, alignment=PP_ALIGN.CENTER)

    tf = add_textbox(s, x + 0.2, 3.2, 2.4, 3.2, "", font_size=12, color=LIGHT_TEXT)
    add_bullets(tf, items, font_size=12, color=LIGHT_TEXT)

add_footer(s, 12)
add_speaker_notes(s, "Mapped to the four judging criteria. The competition brief says: we are not looking for the most profitable strategy; we are looking for the best thinking. That is exactly what we have built.")


# =============================================================================
# SLIDE 13: THANK YOU
# =============================================================================
s = dark_slide()
add_textbox(s, 1.5, 2.0, 10.3, 1.0, "Thank You",
            font_size=48, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 3.2, 10.3, 0.6,
            "ImNuza & Xynerss  |  Artemis Quant Competition 2026  |  Track 1",
            font_size=20, color=GREY, alignment=PP_ALIGN.CENTER)
gold_line(s, 4.5, 4.0, 4.3)
add_textbox(s, 1.5, 4.4, 10.3, 0.5,
            "lindsey@artemisanalytics.xyz",
            font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 5.0, 10.3, 0.5,
            "Code and data: [GitHub repo URL]",
            font_size=16, color=GREY, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 5.5, 10.3, 0.5,
            "Full research report: [PDF link]",
            font_size=16, color=GREY, alignment=PP_ALIGN.CENTER)
add_footer(s, 13)
add_speaker_notes(s, "Thank you. We are happy to take questions. The research report has full details on everything covered today: methodology, data sources, backtest construction, and a complete tested-and-rejected log.")


# =============================================================================
# Save
# =============================================================================
prs.save(str(OUT_PATH))
print(f"Saved: {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
